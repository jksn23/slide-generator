from pptx import Presentation
from zipfile import ZipFile

from core.generator import generate_pptx
from core.models import SlideItem, SlideType, SpeakerLine
from core.template_engine import TemplateResolver


def _shape_texts(presentation):
    return "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))


def test_pptx_renderer_uses_template_and_creates_square_slide(tmp_path):
    output = tmp_path / "rendered.pptx"

    generate_pptx(
        [
            SlideItem(type=SlideType.COVER, content="TATA IBADAH"),
            SlideItem(type=SlideType.SONG_LYRICS, content="Haleluya\nAmin"),
        ],
        str(output),
    )

    presentation = Presentation(str(output))
    assert len(presentation.slides) == 2
    assert presentation.slide_width == presentation.slide_height
    assert "Haleluya" in _shape_texts(presentation)


def test_pptx_renderer_outputs_rich_speaker_text(tmp_path):
    output = tmp_path / "speaker.pptx"

    generate_pptx(
        [
            SlideItem(
                type=SlideType.LITURGY_DIALOG,
                content="Tuhan menyertai saudara.",
                speaker="P",
                speaker_lines=[SpeakerLine("P", "Tuhan menyertai saudara.")],
            )
        ],
        str(output),
    )

    presentation = Presentation(str(output))
    text = _shape_texts(presentation)
    assert "P :" in text
    assert "Tuhan menyertai saudara." in text


def test_pptx_renderer_applies_speaker_line_colors(tmp_path):
    output = tmp_path / "speaker_colors.pptx"

    generate_pptx(
        [
            SlideItem(
                type=SlideType.LITURGY_DIALOG,
                content="P : Tetapi TUHAN\nJ : Amin\nP+J : Amin",
                speaker_lines=[
                    SpeakerLine("P", "Tetapi TUHAN bersemayam untuk selama-lamanya."),
                    SpeakerLine("J", "Amin"),
                    SpeakerLine("P+J", "Amin"),
                ],
            )
        ],
        str(output),
    )

    presentation = Presentation(str(output))
    shape = [shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text][-1]
    paragraphs = shape.text_frame.paragraphs

    assert paragraphs[0].runs[0].font.color.rgb == paragraphs[0].runs[1].font.color.rgb
    assert str(paragraphs[0].runs[0].font.color.rgb) == "FFFFFF"
    assert str(paragraphs[1].runs[0].font.color.rgb) == "F2C94C"
    assert str(paragraphs[1].runs[1].font.color.rgb) == "F2C94C"
    assert str(paragraphs[2].runs[0].font.color.rgb) == "F2C94C"


def test_pptx_renderer_preserves_dialog_font_and_continuation_color(tmp_path):
    output = tmp_path / "speaker_fonts.pptx"

    generate_pptx(
        [
            SlideItem(
                type=SlideType.LITURGY_DIALOG,
                content="P : Tetapi TUHAN\nlanjutan\nJ : Amin\nP+J : Amin",
                speaker_lines=[
                    SpeakerLine("P", "Tetapi TUHAN bersemayam untuk selama-lamanya."),
                    SpeakerLine("", "Ia telah menegakkan takhta-Nya untuk menghakimi."),
                    SpeakerLine("J", "Amin"),
                    SpeakerLine("P+J", "Amin"),
                ],
                metadata={"style": {"font_family": "Arial", "font_size": 41}},
            )
        ],
        str(output),
    )

    presentation = Presentation(str(output))
    shape = [shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text][-1]
    paragraphs = shape.text_frame.paragraphs
    runs = [run for paragraph in paragraphs for run in paragraph.runs if run.text]

    assert all(run.font.name == "Arial" for run in runs)
    assert all(run.font.size.pt == 41 for run in runs)
    assert str(paragraphs[0].runs[-1].font.color.rgb) == "FFFFFF"
    assert str(paragraphs[1].runs[0].font.color.rgb) == "FFFFFF"
    assert str(paragraphs[2].runs[-1].font.color.rgb) == "F2C94C"
    assert str(paragraphs[3].runs[-1].font.color.rgb) == "F2C94C"


def test_generate_pptx_applies_ui_font_overrides(tmp_path):
    output = tmp_path / "font_override.pptx"

    generate_pptx(
        [SlideItem(type=SlideType.COVER, content="TATA IBADAH")],
        str(output),
        font_family="Arial",
        font_sizes={"title": 72, "lyric": 44, "liturgi": 36},
    )

    presentation = Presentation(str(output))
    text_shape = next(shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text)
    run = text_shape.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Arial"
    assert run.font.size.pt == 72


def test_generate_pptx_applies_font_families_by_slide_category(tmp_path):
    output = tmp_path / "font_categories.pptx"

    generate_pptx(
        [
            SlideItem(type=SlideType.SECTION, content="BAGIAN"),
            SlideItem(type=SlideType.SONG_TITLE, content="Menyanyi"),
            SlideItem(type=SlideType.SONG_LYRICS, content="Haleluya"),
            SlideItem(type=SlideType.LITURGY_DIALOG, content="P : Amin", speaker_lines=[SpeakerLine("P", "Amin")]),
        ],
        str(output),
        font_families={
            "heading": "Arial",
            "song_title": "Calibri",
            "lyric": "Tahoma",
            "liturgi": "Verdana",
        },
    )

    presentation = Presentation(str(output))
    fonts = []
    for slide in presentation.slides:
        shape = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text][-1]
        run = next(run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text)
        fonts.append(run.font.name)

    assert fonts == ["Arial", "Calibri", "Tahoma", "Verdana"]


def test_generate_pptx_supports_morph_transition_without_corrupting_file(tmp_path):
    output = tmp_path / "morph.pptx"

    generate_pptx(
        [SlideItem(type=SlideType.COVER, content="TATA IBADAH")],
        str(output),
        transition="morph",
    )

    presentation = Presentation(str(output))
    assert len(presentation.slides) == 1
    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "morph" in slide_xml
    assert "p14" in slide_xml


def test_preview_and_export_read_font_size_from_slide_metadata(tmp_path):
    output = tmp_path / "metadata_font.pptx"
    slide = SlideItem(
        type=SlideType.SONG_LYRICS,
        content="Haleluya",
        metadata={"style": {"font_family": "Arial", "font_size": 57}},
    )

    preview_style = TemplateResolver().resolve(slide)
    generate_pptx([slide], str(output))

    presentation = Presentation(str(output))
    text_shape = next(shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text)
    run = text_shape.text_frame.paragraphs[0].runs[0]
    assert preview_style["font_family"] == "Arial"
    assert preview_style["font_size"] == 57
    assert run.font.name == "Arial"
    assert run.font.size.pt == 57


def test_song_lyrics_all_export_runs_use_default_style(tmp_path):
    output = tmp_path / "song_runs.pptx"
    slide = SlideItem(
        type=SlideType.SONG_LYRICS,
        content="Baris satu\nBaris dua\nBaris tiga",
        metadata={"style": {"font_family": "Arial", "font_size": 49, "color": "#FFFFFF"}},
    )

    generate_pptx([slide], str(output))

    presentation = Presentation(str(output))
    shape = [shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text][-1]
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text]
    assert len(runs) == 3
    assert all(run.font.name == "Arial" for run in runs)
    assert all(run.font.size.pt == 49 for run in runs)
    assert all(str(run.font.color.rgb) == "FFFFFF" for run in runs)


def test_export_text_box_is_centered_on_landscape_slide(tmp_path):
    output = tmp_path / "centered_wide.pptx"

    generate_pptx(
        [SlideItem(type=SlideType.SONG_TITLE, content="Menyanyi KJ No. 4\nHai Mari Sembah")],
        str(output),
        aspect_ratio="landscape_16_9",
    )

    presentation = Presentation(str(output))
    text_shape = next(shape for shape in presentation.slides[0].shapes if hasattr(shape, "text") and shape.text)
    shape_center_x = text_shape.left + text_shape.width // 2
    shape_center_y = text_shape.top + text_shape.height // 2
    slide_center_x = presentation.slide_width // 2
    slide_center_y = presentation.slide_height // 2
    tolerance = presentation.slide_width * 0.01

    assert abs(shape_center_x - slide_center_x) <= tolerance
    assert abs(shape_center_y - slide_center_y) <= tolerance
