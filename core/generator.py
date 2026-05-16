import os
import tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from core.models import SlideType, SlideItem

# PPTX namespace for transition XML
PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Formats natively supported by python-pptx
_PPTX_SUPPORTED = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tiff", ".tif", ".wmf"}


def _ensure_supported_image(image_path: str) -> tuple:
    """Return (path_to_use, temp_file_or_None).
    If the image format is not supported by python-pptx, convert it to PNG
    in a temp file and return that path. Caller must delete temp_file when done.
    """
    ext = os.path.splitext(image_path)[1].lower()
    if ext in _PPTX_SUPPORTED:
        return image_path, None

    try:
        from PIL import Image
        img = Image.open(image_path).convert("RGBA")
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp.name, format="PNG")
        tmp.close()
        return tmp.name, tmp.name
    except Exception as e:
        raise RuntimeError(f"Tidak dapat mengonversi gambar '{image_path}': {e}")


def _apply_transition(slide, transition: str):
    """Inject OOXML transition element into a slide XML tree."""
    if not transition:
        return
    spTree = slide.shapes._spTree
    # Build namespace-qualified tag
    nsmap = {"p": PPTX_NS}
    trans_tag = f"{{{PPTX_NS}}}transition"

    # Remove existing transition if any
    for existing in spTree.getparent().findall(trans_tag):
        spTree.getparent().remove(existing)

    # Build the correct OXML depending on type
    slide_xml = slide._element
    if transition == "fade":
        xml_str = '<p:transition xmlns:p="{ns}" spd="med"><p:fade/></p:transition>'.format(ns=PPTX_NS)
    elif transition == "wipe":
        xml_str = '<p:transition xmlns:p="{ns}" spd="med"><p:wipe dir="l"/></p:transition>'.format(ns=PPTX_NS)
    elif transition == "push":
        xml_str = '<p:transition xmlns:p="{ns}" spd="med"><p:push dir="l"/></p:transition>'.format(ns=PPTX_NS)
    elif transition == "zoom":
        xml_str = '<p:transition xmlns:p="{ns}" spd="med"><p:zoom dir="out"/></p:transition>'.format(ns=PPTX_NS)
    else:
        return

    trans_el = etree.fromstring(xml_str)
    slide_xml.append(trans_el)


def _add_bg_image(slide, image_path: str, prs: Presentation):
    """Add an image as slide background (full-bleed, behind all shapes)."""
    if not image_path or not os.path.exists(image_path):
        return
    # Convert to supported format if necessary (e.g. WebP → PNG)
    use_path, tmp_path = _ensure_supported_image(image_path)
    try:
        left = top = Emu(0)
        width = prs.slide_width
        height = prs.slide_height
        pic = slide.shapes.add_picture(use_path, left, top, width, height)
        # Move the picture to the bottom of the z-order
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)


def generate_pptx(
    slides: list,
    output_path: str,
    font_family: str = "Segoe UI",
    font_sizes: dict = None,
    transition: str = None,
):
    if font_sizes is None:
        font_sizes = {}

    fs_title = font_sizes.get("title", 60)    # Cover, Bagian, Penutup, Judul Lagu
    fs_lyric = font_sizes.get("lyric", 48)    # Lirik Lagu
    fs_lit   = font_sizes.get("liturgi", 40)  # Liturgi, Bacaan, Instruksi

    prs = Presentation()
    prs.slide_width  = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]

    for item in slides:
        if not item.included:
            continue

        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()

        # Default textbox geometry (overridden per type below)
        left, top, width, height = Inches(1), Inches(1), Inches(14), Inches(7)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # ── COVER ───────────────────────────────────────────
        if item.slide_type == SlideType.COVER:
            fill.fore_color.rgb = RGBColor(0x1F, 0x4D, 0x3A)
            p.text = item.content
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.size = Pt(fs_title)
            p.font.bold = True
            p.font.name = font_family

        # ── INSTRUKSI ───────────────────────────────────────
        elif item.slide_type == SlideType.INSTRUKSI:
            fill.fore_color.rgb = RGBColor(0x2F, 0x80, 0xED)
            p.text = item.content
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.size = Pt(fs_lit)
            p.font.name = font_family

        # ── BAGIAN ──────────────────────────────────────────
        elif item.slide_type == SlideType.BAGIAN:
            fill.fore_color.rgb = RGBColor(0x12, 0x33, 0x26)
            txBox.top = Inches(3.5)
            p.text = item.content
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.size = Pt(fs_title)
            p.font.bold = True
            p.font.name = font_family
            # Overlay background image if set
            if item.bg_image:
                _add_bg_image(slide, item.bg_image, prs)

        # ── JUDUL LAGU ──────────────────────────────────────
        elif item.slide_type == SlideType.JUDUL_LAGU:
            fill.fore_color.rgb = RGBColor(0x1F, 0x4D, 0x3A)
            txBox.top = Inches(3.5)
            p.text = item.content
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0xF2, 0xC9, 0x4C)
            p.font.size = Pt(fs_title)
            p.font.bold = True
            p.font.name = font_family

        # ── LIRIK LAGU ──────────────────────────────────────
        elif item.slide_type == SlideType.LIRIK_LAGU:
            fill.fore_color.rgb = RGBColor(0x20, 0x23, 0x22)
            if item.title:
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
                tp = title_box.text_frame.paragraphs[0]
                tp.text = item.title
                tp.alignment = PP_ALIGN.CENTER
                tp.font.color.rgb = RGBColor(0x8A, 0x92, 0x8A)
                tp.font.size = Pt(24)
                tp.font.name = font_family
            txBox.top = Inches(2)
            tf.text = item.content
            for par in tf.paragraphs:
                par.alignment = PP_ALIGN.CENTER
                par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                par.font.size = Pt(fs_lyric)
                par.font.name = font_family

        # ── LITURGI ─────────────────────────────────────────
        elif item.slide_type == SlideType.LITURGI:
            fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xF6)

            # Speaker label
            if item.speaker:
                sp_box = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(14), Inches(0.9))
                sp_p = sp_box.text_frame.paragraphs[0]
                sp_p.text = f"{item.speaker} :"
                sp_p.font.bold = True
                sp_p.font.size = Pt(int(fs_lit * 0.85))
                sp_p.font.name = font_family
                if "J" in item.speaker:
                    sp_p.font.color.rgb = RGBColor(0x8A, 0x6D, 0x1D)
                    sp_p.alignment = PP_ALIGN.RIGHT
                elif "+" in item.speaker:
                    sp_p.font.color.rgb = RGBColor(0x1F, 0x4D, 0x3A)
                    sp_p.alignment = PP_ALIGN.CENTER
                else:  # P
                    sp_p.font.color.rgb = RGBColor(0x1F, 0x4D, 0x3A)
                    sp_p.alignment = PP_ALIGN.LEFT

            txBox.top = Inches(1.5)
            tf.text = item.content
            for par in tf.paragraphs:
                par.font.color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
                par.font.size = Pt(fs_lit)
                par.font.name = font_family

                # Alignment: nyanyian → center; bacaan → left (or speaker-based)
                if item.is_nyanyian:
                    par.alignment = PP_ALIGN.CENTER
                else:
                    if item.speaker == "J":
                        par.alignment = PP_ALIGN.RIGHT
                    elif item.speaker and "+" in item.speaker:
                        par.alignment = PP_ALIGN.CENTER
                    else:
                        par.alignment = PP_ALIGN.LEFT

        # ── BACAAN ──────────────────────────────────────────
        elif item.slide_type == SlideType.BACAAN:
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf.text = item.content
            for par in tf.paragraphs:
                par.alignment = PP_ALIGN.LEFT
                par.font.color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
                par.font.size = Pt(fs_lit)
                par.font.name = font_family

        # ── PENUTUP ─────────────────────────────────────────
        elif item.slide_type == SlideType.PENUTUP:
            fill.fore_color.rgb = RGBColor(0x1F, 0x4D, 0x3A)
            txBox.top = Inches(3.5)
            tf.text = item.content
            for par in tf.paragraphs:
                par.alignment = PP_ALIGN.CENTER
                par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                par.font.size = Pt(fs_title)
                par.font.bold = True
                par.font.name = font_family

        # Apply transition to every slide
        _apply_transition(slide, transition)

    prs.save(output_path)
