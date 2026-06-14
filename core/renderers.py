import os
import tempfile
import uuid
from typing import Iterable

from lxml import etree
from lxml.builder import ElementMaker
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from core.models import SlideDeck, SlideItem, SlideBackground, SlideType
from core.template_engine import TemplateResolver


PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
PPTX_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_PPTX_SUPPORTED = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tiff", ".tif", ".wmf"}


def _rgb(value: str) -> RGBColor:
    value = (value or "#FFFFFF").lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _alignment(value: str) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(value, PP_ALIGN.CENTER)


def _ensure_supported_image(image_path: str) -> tuple[str, str | None]:
    ext = os.path.splitext(image_path)[1].lower()
    if ext in _PPTX_SUPPORTED:
        return image_path, None

    try:
        from PIL import Image

        image = Image.open(image_path).convert("RGBA")
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        image.save(tmp.name, format="PNG")
        tmp.close()
        return tmp.name, tmp.name
    except Exception as exc:
        raise RuntimeError(f"Tidak dapat mengonversi gambar '{image_path}': {exc}") from exc


class BackgroundRenderer:
    def render(self, slide, prs: Presentation, slide_item: SlideItem, style: dict) -> None:
        background = self._resolve_background(slide_item, style)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(background.color or "#000000")

        if background.image and os.path.exists(background.image):
            self._add_image(slide, prs, background.image)

        if background.overlay_opacity > 0:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(0),
                Emu(0),
                prs.slide_width,
                prs.slide_height,
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = _rgb(background.overlay_color)
            overlay.fill.transparency = int(max(0.0, min(background.overlay_opacity, 1.0)) * 100)
            overlay.line.fill.background()

    def _resolve_background(self, slide_item: SlideItem, style: dict) -> SlideBackground:
        style_bg = style.get("background", {})
        background = SlideBackground.from_any(style_bg) or SlideBackground()
        if slide_item.background:
            if slide_item.background.color:
                background.color = slide_item.background.color
            if slide_item.background.image:
                background.image = slide_item.background.image
            background.overlay_color = slide_item.background.overlay_color or background.overlay_color
            background.overlay_opacity = slide_item.background.overlay_opacity
        return background

    def _add_image(self, slide, prs: Presentation, image_path: str) -> None:
        use_path, tmp_path = _ensure_supported_image(image_path)
        try:
            picture = slide.shapes.add_picture(use_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
            sp_tree = slide.shapes._spTree
            sp_tree.remove(picture._element)
            sp_tree.insert(2, picture._element)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except (PermissionError, OSError):
                    pass


class TextRenderer:
    REFERENCE_WIDTH = 10.0
    REFERENCE_HEIGHT = 10.0

    def render(self, slide, prs: Presentation, slide_item: SlideItem, style: dict) -> None:
        left, top, width, height = self._scaled_margin(prs, style.get("margin", {}))
        font_size = int(style.get("font_size", 42))

        if style.get("text_shadow"):
            self._render_text_box(slide, slide_item, style, left + Inches(0.04), top + Inches(0.04), width, height, font_size, "#000000", True)
        self._render_text_box(slide, slide_item, style, left, top, width, height, font_size, style.get("color", "#FFFFFF"), False)

    def _scaled_margin(self, prs: Presentation, margin: dict):
        slide_width = prs.slide_width / Inches(1)
        slide_height = prs.slide_height / Inches(1)
        scale_x = slide_width / self.REFERENCE_WIDTH
        scale_y = slide_height / self.REFERENCE_HEIGHT
        return (
            Inches(float(margin.get("left", 0.8)) * scale_x),
            Inches(float(margin.get("top", 0.8)) * scale_y),
            Inches(float(margin.get("width", 8.4)) * scale_x),
            Inches(float(margin.get("height", 8.4)) * scale_y),
        )

    def _render_text_box(self, slide, slide_item: SlideItem, style: dict, left, top, width, height, font_size: int, color: str, shadow: bool) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE if style.get("vertical_align") == "middle" else MSO_ANCHOR.TOP

        if slide_item.speaker_lines:
            self._render_speaker_lines(frame, slide_item, style, font_size, shadow)
        else:
            self._render_plain_text(frame, slide_item.content, style, font_size, color, shadow)

    def _render_plain_text(self, frame, content: str, style: dict, font_size: int, color: str, shadow: bool) -> None:
        lines = (content or "").splitlines() or [""]
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = _alignment(style.get("align", "center"))
            run = paragraph.add_run()
            run.text = line
            self._apply_font(run, style, font_size, color, shadow)

    def _render_speaker_lines(self, frame, slide_item: SlideItem, style: dict, font_size: int, shadow: bool) -> None:
        speaker_colors = style.get("speaker_colors", {})
        last_speaker = ""
        for index, speaker_line in enumerate(slide_item.speaker_lines):
            speaker = (speaker_line.speaker or "").strip()
            if speaker:
                last_speaker = speaker
            effective_speaker = speaker or last_speaker
            line_color = speaker_colors.get(effective_speaker, style.get("color", "#FFFFFF"))
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = self._speaker_alignment(effective_speaker)
            if speaker:
                speaker_run = paragraph.add_run()
                speaker_run.text = f"{speaker} : "
                self._apply_font(
                    speaker_run,
                    {**style, "bold": True},
                    font_size,
                    line_color,
                    shadow,
                )
            text_run = paragraph.add_run()
            text_run.text = speaker_line.text
            self._apply_font(
                text_run,
                style,
                font_size,
                line_color,
                shadow,
            )

    def _apply_font(self, run, style: dict, font_size: int, color: str, shadow: bool) -> None:
        run.font.name = style.get("font_family", "Segoe UI")
        run.font.size = Pt(font_size)
        run.font.bold = bool(style.get("bold", False))
        run.font.color.rgb = _rgb("#000000" if shadow else color)

    def _speaker_alignment(self, speaker: str) -> PP_ALIGN:
        if speaker == "J":
            return PP_ALIGN.RIGHT
        if "+" in speaker:
            return PP_ALIGN.CENTER
        return PP_ALIGN.LEFT


class PPTXRenderer:
    TRANSITION_XML = {
        "fade": '<p:transition xmlns:p="{ns}" spd="med"><p:fade/></p:transition>',
        "wipe": '<p:transition xmlns:p="{ns}" spd="med"><p:wipe dir="l"/></p:transition>',
        "push": '<p:transition xmlns:p="{ns}" spd="med"><p:push dir="l"/></p:transition>',
        "zoom": '<p:transition xmlns:p="{ns}" spd="med"><p:zoom dir="out"/></p:transition>',
        "morph": (
            '<p:transition xmlns:p="{ns}" xmlns:p14="{p14_ns}" spd="med">'
            '<p:extLst>'
            '<p:ext uri="{{B2D3F22D-AB1C-4A42-A4DB-59C51D907E4D}}">'
            '<p14:morph transition="byObject"/>'
            '</p:ext>'
            '</p:extLst>'
            '</p:transition>'
        ),
    }

    def __init__(self, template_name: str = "gmim_default") -> None:
        self.resolver = TemplateResolver(template_name)
        self.background_renderer = BackgroundRenderer()
        self.text_renderer = TextRenderer()

    def render(
        self,
        slides: SlideDeck | Iterable[SlideItem],
        output_path: str,
        aspect_ratio: str = "square",
        transition: str | None = None,
        custom_breakpoints: dict[int, str] | None = None,
    ) -> None:
        deck_slides = slides.slides if isinstance(slides, SlideDeck) else list(slides)
        ratio = self.resolver.aspect_ratio(aspect_ratio)
        prs = Presentation()
        prs.slide_width = Inches(ratio["width"])
        prs.slide_height = Inches(ratio["height"])
        blank_layout = prs.slide_layouts[6]

        slide_sections_mapping = []
        
        # Variabel pelacak untuk custom section
        active_custom_section = "Tata Ibadah"
        
        # Variabel pelacak untuk pewarisan background
        last_sec_name = None
        active_section_bg = None

        for index, slide_item in enumerate(deck_slides):
            if not slide_item.include:
                continue
            try:
                slide = prs.slides.add_slide(blank_layout)
                
                # --- LOGIKA PENENTUAN SECTION ---
                if custom_breakpoints is not None:
                    # Jika user memberikan custom breakpoints, cek apakah index saat ini adalah titik potong baru
                    if index in custom_breakpoints:
                        active_custom_section = custom_breakpoints[index]
                    sec_name = active_custom_section
                else:
                    # Fallback ke perilaku lama jika tidak ada custom breakpoints
                    sec_name = slide_item.section if slide_item.section else "Tata Ibadah"
                # ---------------------------------

                slide_sections_mapping.append((sec_name, slide.slide_id))

                # --- LOGIKA PEWARISAN BACKGROUND SECTION ---
                # 1. Deteksi jika berpindah ke section baru
                if sec_name != last_sec_name:
                    active_section_bg = None
                    last_sec_name = sec_name
                    
                # 2. Update background aktif jika slide ini memiliki pengaturan spesifik
                if slide_item.background and (slide_item.background.image or slide_item.background.color):
                    active_section_bg = slide_item.background
                
                # 3. Wariskan background jika slide ini kosong tapi ada background aktif di section ini
                elif active_section_bg is not None:
                    if slide_item.background is None:
                        slide_item.background = SlideBackground()
                    
                    # Salin properti dari background section yang aktif
                    slide_item.background.image = active_section_bg.image
                    slide_item.background.color = active_section_bg.color
                    slide_item.background.overlay_color = active_section_bg.overlay_color
                    slide_item.background.overlay_opacity = active_section_bg.overlay_opacity
                # -------------------------------------------

                style = self.resolver.resolve(slide_item)
                
                if getattr(slide_item, 'is_absolute_layout', False):
                    slide_w_inches = prs.slide_width
                    slide_h_inches = prs.slide_height
                    
                    left = slide_w_inches * slide_item.title_pos_x
                    top = slide_h_inches * slide_item.title_pos_y
                    width = slide_w_inches * slide_item.title_width
                    height = slide_h_inches * slide_item.title_height
                    
                    if slide_item.background and (slide_item.background.image or slide_item.background.color):
                        self.background_renderer.render(slide, prs, slide_item, style)
                        
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    lines = slide_item.content.split('\n')
                    for i, line in enumerate(lines):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = line
                        if slide_item.title_font_profile:
                            p.font.name = slide_item.title_font_profile.get("family", "Arial")
                            p.font.size = Pt(slide_item.title_font_profile.get("size", 44))
                            if slide_item.title_font_profile.get("color"):
                                p.font.color.rgb = _rgb(slide_item.title_font_profile["color"])
                            if slide_item.title_font_profile.get("bold"):
                                p.font.bold = True
                else:
                    self.background_renderer.render(slide, prs, slide_item, style)
                    if slide_item.type != SlideType.BLANK:
                        self.text_renderer.render(slide, prs, slide_item, style)
                        
                self._apply_transition(slide, transition)
            except Exception as e:
                import logging
                logging.getLogger(__name__).warning(f"Gagal merender slide {slide_item.id}: {e}")
                continue

        # Injeksi XML Sections sebelum file disimpan
        self._inject_sections(prs, slide_sections_mapping)

        prs.save(output_path)

    def _apply_transition(self, slide, transition: str | None) -> None:
        if not transition:
            return
        slide_xml = slide._element
        tag = f"{{{PPTX_NS}}}transition"
        for existing in slide_xml.findall(tag):
            slide_xml.remove(existing)
        transition_xml = self.TRANSITION_XML.get(transition, self.TRANSITION_XML["fade"])
        transition_element = etree.fromstring(transition_xml.format(ns=PPTX_NS, p14_ns=PPTX_P14_NS))
        
        cSld = slide_xml.find(f"{{{PPTX_NS}}}cSld")
        if cSld is not None:
            index = slide_xml.index(cSld)
            slide_xml.insert(index + 1, transition_element)
        else:
            slide_xml.append(transition_element)

    def _inject_sections(self, prs, mapping):
        if not mapping:
            return

        P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
        P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
        EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

        # 1. Kelompokkan ID Slide ke dalam blok yang berurutan (Contiguous Blocks)
        blocks = []
        current_section = None
        current_ids = []

        for sec_name, sld_id in mapping:
            if sec_name != current_section:
                if current_section is not None:
                    blocks.append((current_section, current_ids))
                current_section = sec_name
                current_ids = [sld_id]
            else:
                current_ids.append(sld_id)
        if current_section is not None:
            blocks.append((current_section, current_ids))

        # 2. Akses root XML
        presentation_xml = prs.element
        extLst = presentation_xml.find(f'./{{{P_NS}}}extLst')
        
        # 3. Buat <p:extLst> jika belum ada
        if extLst is None:
            extLst = etree.Element(f'{{{P_NS}}}extLst')
            presentation_xml.append(extLst)

        # 4. Buat sub-elemen ekstensi khusus
        section_ext = etree.Element(f'{{{P_NS}}}ext', uri=EXT_URI)
        extLst.append(section_ext)

        builder = ElementMaker(namespace=P14_NS, nsmap={'p14': P14_NS})
        sectionList = builder.sectionLst()
        section_ext.append(sectionList)

        # 5. Injeksi XML
        for section_name, slide_ids in blocks:
            sec_id = f"{{{str(uuid.uuid4()).upper()}}}"
            section_node = etree.Element(f'{{{P14_NS}}}section', name=str(section_name), id=sec_id)
            slideIdLst = etree.Element(f'{{{P14_NS}}}sldIdLst')
            
            for sld_id in slide_ids:
                slideId = etree.Element(f'{{{P14_NS}}}sldId', id=str(sld_id))
                slideIdLst.append(slideId)
            
            section_node.append(slideIdLst)
            sectionList.append(section_node)
